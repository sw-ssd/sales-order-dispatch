#!/usr/bin/env python3
"""Generate a customer-facing PPT report for 多公司訂出貨系統 1.0."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

OUTPUT = os.path.join(os.path.dirname(__file__), "多公司訂出貨系統_1.0_客戶報告.pptx")
FONT = "PingFang TC"
ACCENT = RGBColor(0x1E, 0x4D, 0x8C)  # deep blue
ACCENT_LIGHT = RGBColor(0x3A, 0x7B, 0xD5)  # lighter blue
ACCENT_DARK = RGBColor(0x12, 0x35, 0x60)  # darker blue
DARK = RGBColor(0x22, 0x22, 0x22)
LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
SOFT_BG = RGBColor(0xF5, 0xF8, 0xFC)  # very light blue-gray


def set_text(run, text, size=18, bold=False, color=DARK):
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def apply_gradient(shape, color1, color2, angle=270):
    """Apply a two-stop linear gradient fill to a shape."""
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = angle
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color2
    fill.gradient_stops[1].position = 1.0
    shape.line.fill.background()


def add_decorative_circle(slide, prs, left, top, width, height, color, transparency=0.85):
    """Add a semi-transparent circle/oval for visual decoration."""
    shape = slide.shapes.add_shape(9, left, top, width, height)  # 9 = oval
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # python-pptx does not expose transparency directly; use a lighter color instead
    return shape


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # gradient background
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    apply_gradient(background, ACCENT_DARK, ACCENT_LIGHT, angle=315)

    # decorative circles
    add_decorative_circle(slide, prs, Inches(9.5), Inches(-1.5), Inches(5), Inches(5), ACCENT_LIGHT)
    add_decorative_circle(slide, prs, Inches(-2), Inches(4.5), Inches(4.5), Inches(4.5), ACCENT)

    # title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.4), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_text(p.add_run(), title, size=48, bold=True, color=LIGHT)

    # subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.4), Inches(1))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_text(p.add_run(), subtitle, size=24, color=LIGHT)

    # footer
    footer = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.4), Inches(0.5))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_text(p.add_run(), "修訂號 v1.0.2 · 2026-07-17", size=14, color=LIGHT)
    return slide


def add_section_slide(prs, number, title):
    """Big section divider."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    apply_gradient(background, ACCENT, ACCENT_DARK, angle=270)

    # large decorative number (translucent effect via lighter color)
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.4), Inches(2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_text(p.add_run(), number, size=120, bold=True, color=ACCENT_LIGHT)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.4), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_text(p.add_run(), title, size=42, bold=True, color=LIGHT)
    return slide


def add_content_slide(prs, title, bullets, note=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # soft page background
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SOFT_BG
    bg.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)

    # left accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), prs.slide_height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    # title bar with gradient
    bar = slide.shapes.add_shape(1, Inches(0.18), Inches(0), Inches(12.82), Inches(1.15))
    apply_gradient(bar, ACCENT, ACCENT_LIGHT, angle=0)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(11.8), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    set_text(p.add_run(), title, size=28, bold=True, color=LIGHT)

    # content card
    card = slide.shapes.add_shape(1, Inches(0.4), Inches(1.35), Inches(12.2), Inches(5.85))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = RGBColor(0xDD, 0xE4, 0xEE)
    card.line.width = Pt(1)
    # send card behind text but above background
    spTree = slide.shapes._spTree
    sp = card._element
    spTree.remove(sp)
    spTree.insert(3, sp)

    # content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(11.6), Inches(5.45))
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.1)
    for idx, bullet in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(16)
        p.text = f"• {bullet}"
        for run in p.runs:
            run.font.name = FONT
            run.font.size = Pt(20)
            run.font.color.rgb = DARK

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.75), Inches(11.6), Inches(0.5))
        tf = note_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        set_text(p.add_run(), note, size=14, color=GRAY)

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, "多公司訂出貨系統 1.0", "下一代訂出貨管理系統規劃報告")

    add_section_slide(prs, "01", "專案目標與價值")
    add_content_slide(prs, "為什麼需要新系統？", [
        "完全擺脫外部系統依賴，客戶、商品、業務主檔改為自建管理。",
        "建立『公司 → 部門』兩層組織架構，支援多公司、多部門獨立運作。",
        "強化倉儲、派車與單據列印流程，讓出貨作業更順暢。",
        "保留現有使用習慣，同時提升權限管控、通知推播與稽核能力。",
    ])

    add_section_slide(prs, "02", "技術架構")
    add_content_slide(prs, "全新 monorepo，統一技術棧", [
        "後端：Go 1.25 + Chi REST + Connect-RPC + Ent ORM + PostgreSQL + Valkey。",
        "Web 中台：SolidJS 1.9 + TypeScript 5.9 + TanStack Query / Router。",
        "行動 App：Flutter 3.35.2 + connectrpc + Sembast 本地快取。",
        "權限：後端 Casbin RBAC with domain，前端 CASL.js。",
        "單據列印：Gotenberg HTML → PDF，支援繁體中文。",
        "Monorepo：pnpm workspace + Turborepo 統一管理全倉任務。",
    ], note="採用全新 monorepo 從頭開發，沿用現有使用習慣但技術架構全面升級。")

    add_section_slide(prs, "03", "使用者與權限")
    add_content_slide(prs, "六種角色，各司其職", [
        "系統管理員（super）：建立公司與部門，管理所有使用者與資料。",
        "公司管理員（company_admin）：單一公司內的超級管理員，可跨部門。",
        "部門管理員（dept_admin）：管理所屬部門的客戶、商品、訂單、派車與人員。",
        "業務（staff）：負責日常下單、客戶與商品維護。",
        "客戶（customer）：外部店家，只能查看自己的訂單與專屬商品。",
        "訪客（guest）：首次 OAuth 登入後自動建立，待管理員審核後才能操作。",
    ], note="後端使用 Casbin、前端使用 CASL.js，實現公司級與部門級雙重權限控管。")

    add_section_slide(prs, "04", "公司識別與公開資訊")
    add_content_slide(prs, "每間公司擁有專屬對外形象", [
        "公司識別標識：可上傳 Logo、設定主色碼（primary color）、公司簡稱/英文名稱。",
        "公開資訊：公司電話、地址、Email、統一編號、服務條款與隱私權政策連結。",
        "動態呈現：登入頁、Web 側邊欄、App 首頁與關於我們頁面會依所屬公司自動切換。",
        "單據表頭：PDF 出貨單據自動帶入公司 Logo 與公開資訊，強化專業形象。",
        "權限控管：僅 super 與 company_admin 可編輯所屬公司識別與公開資訊。",
    ], note="讓同一系統下的不同公司/品牌都能擁有獨立對外識別。")

    add_section_slide(prs, "05", "核心功能總覽")
    add_content_slide(prs, "七大功能模組", [
        "組織與帳號：公司、部門、使用者管理與 OAuth2 / 客戶帳號密碼登入。",
        "主檔管理：客戶總表、商品總表、客戶專屬商品清單、字典檔。",
        "銷售訂單：業務/客戶下單、單位換算、訂單狀態追蹤。",
        "派車規劃：Web 看板拖放排車、即時同步、依出貨日篩選。",
        "單據列印：單車總表、對點單、揀貨單、加工單四種 PDF。",
        "通知推播：Email 通知、App FCM 推播、Web 通知中心。",
        "稽核日誌：登入、下單、修改主檔、列印等關鍵操作全程記錄。",
    ])

    add_section_slide(prs, "06", "登入與安全")
    add_content_slide(prs, "兩種登入方式，彈性又安全", [
        "員工 / 管理員：透過 Google Workspace 帳號 OAuth2 登入。",
        "客戶：使用系統配發的 customer_code + 密碼登入，首次登入強制改密碼。",
        "QR Code 深層連結：業務拜訪時可掃碼帶入帳號，快速登入。",
        "首次 OAuth 登入若為陌生帳號，系統自動設為 guest，待審核後指派部門與角色。",
        "帳號僅可停用，不可刪除，確保資料稽核軌跡完整。",
    ])

    add_section_slide(prs, "07", "訂單流程")
    add_content_slide(prs, "從下單到出貨，一氣呵成", [
        "選擇客戶（業務）或自動帶入自己（客戶）。",
        "帶出客戶專屬商品清單，也可從商品總表手打品名並儲存為別名。",
        "輸入數量、單位、單價、分切規格與特殊切法備註，系統自動換算。",
        "選擇預計出貨日，提交後自動發送 Email 與 App 推播通知。",
        "訂單狀態：待處理 → 處理中 → 已完成；待處理階段可編輯與取消。",
    ])

    add_section_slide(prs, "08", "派車與單據")
    add_content_slide(prs, "派車看板 + 四種單據", [
        "Web 派車看板：以車次為欄、訂單為卡，拖放即可調整派車。",
        "即時同步：多部門同事同時操作時，看板自動更新。",
        "單車總表：彙整各車次所有店家與明細摘要。",
        "對點單：每個店家一張 A4，列出品名、數量、單位（不含價格）。",
        "揀貨單：依車次、倉別、商品分類排序，方便倉庫揀貨。",
        "加工單：區分「加工室揀」與「配送揩」，標示原始與加工後數量。",
    ], note="所有單據皆以 HTML → PDF 產生，A4 輸出，字型支援繁體中文。")

    add_section_slide(prs, "09", "Web 中台介面")
    add_content_slide(prs, "網頁後台主要頁面", [
        "Dashboard：今日待出貨、待處理訂單、快速連結。",
        "公司 / 部門管理：super 專用，建立與維護組織。",
        "公司識別設定：上傳 Logo、設定主色與公開資訊，打造專屬對外形象。",
        "管理人員名單：使用者 CRUD、角色指派、停用、強制登出。",
        "客戶總表、商品總表、客戶專屬商品：部門主檔維護。",
        "訂單管理：列表、新增、編輯、檢視。",
        "派車規劃：Kanban 看板、拖放、列印。",
        "公告管理、通知中心、系統設定。",
    ])

    add_section_slide(prs, "10", "App 行動應用")
    add_content_slide(prs, "業務與客戶都能用手機下單", [
        "身分選擇：登入時區分「我是店家」或「我是業務」。",
        "首頁：公告輪播、最新消息、快速下單入口，並依公司顯示專屬 Logo 與主色。",
        "底部導覽：首頁 / 商品 / 訂單歷史 / 功能。",
        "業務：選客戶 → 帶出專屬商品 → 下單；可管理客戶專屬商品別名。",
        "客戶：只能從自己的專屬商品清單下單，並查看訂單歷史。",
        "關於我們：顯示所屬公司的公開資訊與聯絡方式。",
        "離線快取：登入資訊、客戶、商品、訂單歷史可快取，下拉即可刷新。",
    ])

    add_section_slide(prs, "11", "通知與稽核")
    add_content_slide(prs, "重要事件不漏接，操作全程有記錄", [
        "Email 通知：訂單建立時自動寄送給相關人員。",
        "App 推播：透過 Firebase Cloud Messaging 發送訂單狀態與派車通知。",
        "Web 通知中心：系統通知集中顯示。",
        "稽核日誌：記錄登入、下單、主檔修改、刪除、列印、強制登出、角色變更。",
        "每筆日誌包含操作人、時間、操作類型與異動前後摘要。",
    ])

    add_section_slide(prs, "12", "安全與維運")
    add_content_slide(prs, "資料保護與系統維運", [
        "個人資料保護：客戶與員工資料加密儲存與傳輸，嚴格依角色與部門隔離。",
        "稽核日誌：登入、下單、修改主檔、刪除、列印等關鍵操作全程記錄。",
        "備份策略：資料庫每日備份 + 即時歸檔，備份存放於 Google Cloud Storage。",
        "災難復原：RTO 4 小時 / RPO 1 小時，每半年執行復原演練。",
        "監控告警：基礎設施與業務指標監控，異常即時通知。",
    ], note="1.0 採 Big Bang 切換：測試完成後直接全面上線，不與舊系統並行。")

    add_section_slide(prs, "13", "第一版範圍")
    add_content_slide(prs, "第一版實作重點", [
        "公司、部門、使用者管理與 OAuth2 / 客戶帳號密碼登入。",
        "多公司/部門權限控管、資料庫 RLS 租戶隔離、稽核日誌。",
        "客戶主檔（含地址簿、聯絡人）、商品主檔、客戶專屬商品清單。",
        "倉庫、車次、分切規格、商品分類等部門級字典。",
        "銷售訂單建立、編輯、狀態追蹤、金額計算與通知。",
        "Web 派車看板即時同步與四種出貨單據 PDF 列印。",
        "Flutter App：業務/客戶快速下單、訂單歷史、推播。",
        "公司識別與公開資訊：Logo、主色、關於我們、單據表頭。",
    ])
    add_content_slide(prs, "第一版暫不處理", [
        "供應商退貨授權（vendor return authorization）。",
        "獨立報價單功能：以客戶專屬商品清單取代報價帶入。",
        "簽核流程（approval workflow）。",
        "進銷存即時庫存管理：僅記錄庫別與分庫屬性。",
        "退貨流程。",
        "從現有 sales-order 系統匯入歷史資料：1.0 重新建檔。",
        "AI Agent 整合：預留 capabilities / public_info，不實作完整協定。",
    ], note="這些項目可在後續版本評估擴充。")

    add_section_slide(prs, "14", "預期效益")
    add_content_slide(prs, "導入後可期待的改善", [
        "降低外部系統授權與客製成本，資料主導權回到企業手中。",
        "多公司/部門權限清晰，資料隔離安全，管理更有彈性。",
        "派車與單據流程數位化，減少人工抄寫與溝通錯誤。",
        "業務與客戶都能隨時隨地下單，提升訂單處理效率。",
        "完整稽核日誌與通知機制，強化內控與客戶服務。",
    ])

    add_section_slide(prs, "15", "下一步")
    add_content_slide(prs, "建議推進方式", [
        "確認本報告所列功能範圍與優先順序。",
        "進入 Phase 1 實作：權限/認證 → 主檔 → 訂單 → 派車/列印 → App。",
        "每個階段完成後進行驗收測試，確保符合需求再繼續。",
        "上線前由管理員預先建立員工帳號，避免 OAuth 首次登入審核瓶頸。",
    ])

    # closing
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    apply_gradient(background, ACCENT_DARK, ACCENT_LIGHT, angle=45)

    add_decorative_circle(slide, prs, Inches(10), Inches(-2), Inches(5.5), Inches(5.5), ACCENT_LIGHT)
    add_decorative_circle(slide, prs, Inches(-1.5), Inches(5), Inches(4), Inches(4), ACCENT)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.4), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_text(p.add_run(), "感謝聆聽", size=52, bold=True, color=LIGHT)

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.4), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_text(p.add_run(), "多公司訂出貨系統 1.0", size=24, color=LIGHT)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
